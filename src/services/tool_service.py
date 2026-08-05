from typing import Any, Dict, Optional, List
import subprocess
import os
from datetime import datetime

from src.config.app_config import settings

from src.utils.logger import get_logger

logger = get_logger(__name__)


class ToolService:
    """工具服务 - AI可调用的统一工具接口"""
    
    _instance = None
    
    def __new__(cls, *args, **kwargs):
        if not cls._instance:
            cls._instance = super().__new__(cls)
        return cls._instance
    
    def __init__(self):
        if hasattr(self, '_initialized') and self._initialized:
            return
        
        self._initialized = True
    
    def run_command(self, command: str, timeout: int = 30) -> Dict[str, Any]:
        """执行命令行命令"""
        try:
            result = subprocess.run(
                command,
                shell=True,
                capture_output=True,
                text=True,
                timeout=timeout,
                encoding='utf-8',
                errors='replace'
            )
            return {
                "status": "success" if result.returncode == 0 else "error",
                "stdout": result.stdout,
                "stderr": result.stderr,
                "returncode": result.returncode
            }
        except subprocess.TimeoutExpired:
            return {"status": "error", "message": "命令执行超时"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def ffmpeg_convert(self, input_path: str, output_path: str, options: str = "") -> Dict[str, Any]:
        """使用ffmpeg转换媒体文件"""
        try:
            import ffmpeg
            stream = ffmpeg.input(input_path)
            stream = ffmpeg.output(stream, output_path, **self._parse_ffmpeg_options(options))
            ffmpeg.run(stream, capture_stdout=True, capture_stderr=True)
            return {"status": "success", "output_path": output_path}
        except ImportError:
            # 回退到命令行方式
            cmd = f"ffmpeg -i \"{input_path}\" {options} \"{output_path}\""
            return self.run_command(cmd)
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def _parse_ffmpeg_options(self, options_str: str) -> Dict[str, Any]:
        """解析ffmpeg选项字符串"""
        options = {}
        if options_str:
            parts = options_str.split()
            i = 0
            while i < len(parts):
                if parts[i].startswith('-'):
                    key = parts[i][1:]
                    if i + 1 < len(parts) and not parts[i+1].startswith('-'):
                        options[key] = parts[i+1]
                        i += 2
                    else:
                        options[key] = None
                        i += 1
                else:
                    i += 1
        return options
    
    def audio_transcribe(self, audio_path: str) -> Dict[str, Any]:
        """音频转文字"""
        try:
            import speech_recognition as sr
            r = sr.Recognizer()
            with sr.AudioFile(audio_path) as source:
                audio = r.record(source)
            text = r.recognize_google(audio, language='zh-CN')
            return {"status": "success", "transcription": text}
        except ImportError:
            return {"status": "error", "message": "speech_recognition未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def audio_extract_features(self, audio_path: str) -> Dict[str, Any]:
        """提取音频特征"""
        try:
            import librosa
            y, sr = librosa.load(audio_path)
            tempo, _ = librosa.beat.beat_track(y=y, sr=sr)
            mfcc = librosa.feature.mfcc(y=y, sr=sr)
            
            return {
                "status": "success",
                "tempo": float(tempo),
                "mfcc_mean": float(mfcc.mean()),
                "mfcc_std": float(mfcc.std()),
                "duration": float(len(y) / sr),
                "sample_rate": sr
            }
        except ImportError:
            return {"status": "error", "message": "librosa未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def video_cut(self, video_path: str, start: float, end: float) -> Dict[str, Any]:
        """剪辑视频"""
        try:
            from moviepy.editor import VideoFileClip
            clip = VideoFileClip(video_path)
            subclip = clip.subclip(start, end)
            output_path = video_path.replace('.mp4', '_cut.mp4')
            subclip.write_videofile(output_path)
            
            return {"status": "success", "output_path": output_path}
        except ImportError:
            return {"status": "error", "message": "moviepy未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def video_extract_audio(self, video_path: str) -> Dict[str, Any]:
        """从视频提取音频"""
        try:
            from moviepy.editor import VideoFileClip
            video = VideoFileClip(video_path)
            audio = video.audio
            output_path = video_path.replace('.mp4', '.mp3')
            audio.write_audiofile(output_path)
            
            return {"status": "success", "output_path": output_path}
        except ImportError:
            return {"status": "error", "message": "moviepy未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def audio_convert(self, input_path: str, output_format: str = "mp3") -> Dict[str, Any]:
        """转换音频格式"""
        try:
            from pydub import AudioSegment
            audio = AudioSegment.from_file(input_path)
            output_path = input_path.rsplit('.', 1)[0] + f'.{output_format}'
            audio.export(output_path, format=output_format)
            
            return {"status": "success", "output_path": output_path}
        except ImportError:
            return {"status": "error", "message": "pydub未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def play_audio(self, audio_path: str) -> Dict[str, Any]:
        """播放音频"""
        try:
            import sounddevice as sd
            import soundfile as sf
            
            data, fs = sf.read(audio_path)
            sd.play(data, fs)
            sd.wait()
            
            return {"status": "success", "message": "播放完成"}
        except ImportError:
            return {"status": "error", "message": "sounddevice或soundfile未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def video_info(self, video_path: str) -> Dict[str, Any]:
        """获取视频信息"""
        try:
            import decord
            vr = decord.VideoReader(video_path)
            frame_count = len(vr)
            fps = vr.get_avg_fps()
            
            return {
                "status": "success",
                "frame_count": frame_count,
                "fps": float(fps),
                "duration": float(frame_count / fps),
                "width": vr[0].shape[1],
                "height": vr[0].shape[0]
            }
        except ImportError:
            return {"status": "error", "message": "decord未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def audio_add_effect(self, audio_path: str, effect: str = "reverb") -> Dict[str, Any]:
        """添加音频效果"""
        try:
            import pedalboard
            from pedalboard import Pedalboard, Reverb, Delay, Compressor
            
            board = Pedalboard()
            
            if effect == "reverb":
                board.append(Reverb())
            elif effect == "delay":
                board.append(Delay())
            elif effect == "compressor":
                board.append(Compressor())
            elif effect == "all":
                board.append(Compressor())
                board.append(Reverb())
                board.append(Delay())
            else:
                return {"status": "error", "message": f"未知效果: {effect}"}
            
            audio, sample_rate = pedalboard.io.load(audio_path)
            effected = board(audio, sample_rate)
            
            output_path = audio_path.rsplit('.', 1)[0] + f'_{effect}.wav'
            pedalboard.io.write(output_path, effected, sample_rate)
            
            return {"status": "success", "output_path": output_path}
        except ImportError:
            return {"status": "error", "message": "pedalboard未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def read_docx(self, file_path: str) -> Dict[str, Any]:
        """读取Word文档"""
        try:
            from docx import Document
            doc = Document(file_path)
            text = '\n'.join(para.text for para in doc.paragraphs)
            
            return {"status": "success", "text": text}
        except ImportError:
            return {"status": "error", "message": "python-docx未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def write_docx(self, text: str, output_path: str) -> Dict[str, Any]:
        """写入Word文档"""
        try:
            from docx import Document
            doc = Document()
            doc.add_paragraph(text)
            doc.save(output_path)
            
            return {"status": "success", "output_path": output_path}
        except ImportError:
            return {"status": "error", "message": "python-docx未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def read_excel(self, file_path: str, sheet_name: Optional[str] = None) -> Dict[str, Any]:
        """读取Excel文件"""
        try:
            import openpyxl
            wb = openpyxl.load_workbook(file_path)
            
            if sheet_name:
                ws = wb[sheet_name]
            else:
                ws = wb.active
            
            data = []
            for row in ws.iter_rows(values_only=True):
                data.append(list(row))
            
            return {"status": "success", "data": data, "sheet_names": wb.sheetnames}
        except ImportError:
            return {"status": "error", "message": "openpyxl未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def write_excel(self, data: List[List[Any]], output_path: str) -> Dict[str, Any]:
        """写入Excel文件"""
        try:
            import openpyxl
            wb = openpyxl.Workbook()
            ws = wb.active
            
            for row in data:
                ws.append(row)
            
            wb.save(output_path)
            return {"status": "success", "output_path": output_path}
        except ImportError:
            return {"status": "error", "message": "openpyxl未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def read_pdf(self, file_path: str) -> Dict[str, Any]:
        """读取PDF文件"""
        try:
            import pdfplumber
            with pdfplumber.open(file_path) as pdf:
                text = '\n'.join(page.extract_text() for page in pdf.pages)
            
            return {"status": "success", "text": text[:5000]}
        except ImportError:
            return {"status": "error", "message": "pdfplumber未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def extract_pdf_images(self, file_path: str, output_dir: str = "pdf_images") -> Dict[str, Any]:
        """提取PDF中的图片"""
        try:
            import fitz  # PyMuPDF
            os.makedirs(output_dir, exist_ok=True)
            
            doc = fitz.open(file_path)
            image_count = 0
            
            for page in doc:
                for img in page.get_images(full=True):
                    xref = img[0]
                    base_image = doc.extract_image(xref)
                    image_bytes = base_image["image"]
                    
                    image_path = os.path.join(output_dir, f"image_{image_count}.png")
                    with open(image_path, "wb") as f:
                        f.write(image_bytes)
                    image_count += 1
            
            return {"status": "success", "image_count": image_count, "output_dir": output_dir}
        except ImportError:
            return {"status": "error", "message": "PyMuPDF未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def create_pptx(self, slides: List[Dict[str, str]], output_path: str) -> Dict[str, Any]:
        """创建PPT"""
        try:
            from pptx import Presentation
            prs = Presentation()
            
            for slide_data in slides:
                slide = prs.slides.add_slide(prs.slide_layouts[1])
                if 'title' in slide_data:
                    slide.shapes.title.text = slide_data['title']
                if 'content' in slide_data:
                    slide.placeholders[1].text = slide_data['content']
            
            prs.save(output_path)
            return {"status": "success", "output_path": output_path}
        except ImportError:
            return {"status": "error", "message": "python-pptx未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def create_wordcloud(self, text: str, output_path: str, max_words: int = 100) -> Dict[str, Any]:
        """生成词云"""
        try:
            from wordcloud import WordCloud
            import matplotlib.pyplot as plt
            
            wc = WordCloud(
                font_path="simhei.ttf",
                max_words=max_words,
                width=800,
                height=600,
                background_color="white"
            )
            
            wc.generate(text)
            wc.to_file(output_path)
            
            return {"status": "success", "output_path": output_path}
        except ImportError:
            return {"status": "error", "message": "wordcloud或matplotlib未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def plot_chart(self, data: Dict[str, List[float]], chart_type: str = "line", output_path: str = "chart.png") -> Dict[str, Any]:
        """绘制图表"""
        try:
            import matplotlib.pyplot as plt
            import seaborn as sns
            
            plt.figure(figsize=(10, 6))
            
            if chart_type == "line":
                for label, values in data.items():
                    plt.plot(values, label=label)
            elif chart_type == "bar":
                plt.bar(data.keys(), data.values())
            elif chart_type == "scatter":
                plt.scatter(data.get('x', []), data.get('y', []))
            elif chart_type == "hist":
                plt.hist(data.get('values', []), bins=20)
            else:
                return {"status": "error", "message": f"未知图表类型: {chart_type}"}
            
            plt.legend()
            plt.savefig(output_path)
            plt.close()
            
            return {"status": "success", "output_path": output_path}
        except ImportError:
            return {"status": "error", "message": "matplotlib或seaborn未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def text_segment(self, text: str) -> Dict[str, Any]:
        """中文分词"""
        try:
            import thulac
            thu = thulac.thulac()
            result = thu.cut(text, text=True)
            
            return {"status": "success", "segments": result}
        except ImportError:
            return {"status": "error", "message": "thulac未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def sentiment_analysis(self, text: str) -> Dict[str, Any]:
        """情感分析"""
        try:
            from snownlp import SnowNLP
            s = SnowNLP(text)
            
            return {
                "status": "success",
                "sentiment": float(s.sentiments),
                "keywords": s.keywords(5)
            }
        except ImportError:
            return {"status": "error", "message": "snownlp未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def text_summary(self, text: str, ratio: float = 0.2) -> Dict[str, Any]:
        """文本摘要"""
        try:
            from gensim.summarization import summarize
            summary = summarize(text, ratio=ratio)
            
            return {"status": "success", "summary": summary}
        except ImportError:
            return {"status": "error", "message": "gensim未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def generate_json(self, data: Dict[str, Any], output_path: str) -> Dict[str, Any]:
        """生成JSON文件"""
        try:
            import orjson
            with open(output_path, 'wb') as f:
                f.write(orjson.dumps(data, option=orjson.OPT_INDENT_2))
            
            return {"status": "success", "output_path": output_path}
        except ImportError:
            # 回退到标准json
            import json
            with open(output_path, 'w', encoding='utf-8') as f:
                json.dump(data, f, indent=2, ensure_ascii=False)
            return {"status": "success", "output_path": output_path}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def read_json(self, file_path: str) -> Dict[str, Any]:
        """读取JSON文件"""
        try:
            import orjson
            with open(file_path, 'rb') as f:
                data = orjson.loads(f.read())
            
            return {"status": "success", "data": data}
        except ImportError:
            import json
            with open(file_path, 'r', encoding='utf-8') as f:
                data = json.load(f)
            return {"status": "success", "data": data}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def zip_files(self, files: List[str], output_path: str) -> Dict[str, Any]:
        """压缩文件"""
        try:
            import zipfile
            with zipfile.ZipFile(output_path, 'w', zipfile.ZIP_DEFLATED) as zf:
                for file in files:
                    if os.path.exists(file):
                        zf.write(file)
            
            return {"status": "success", "output_path": output_path}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def unzip_file(self, zip_path: str, output_dir: str) -> Dict[str, Any]:
        """解压文件"""
        try:
            import zipfile
            with zipfile.ZipFile(zip_path, 'r') as zf:
                zf.extractall(output_dir)
            
            return {"status": "success", "output_dir": output_dir}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def get_current_time(self) -> Dict[str, Any]:
        """获取当前时间"""
        return {
            "status": "success",
            "datetime": datetime.now().isoformat(),
            "date": datetime.now().strftime('%Y-%m-%d'),
            "time": datetime.now().strftime('%H:%M:%S')
        }
    
    def get_system_info(self) -> Dict[str, Any]:
        """获取系统信息"""
        try:
            import psutil
            return {
                "status": "success",
                "cpu_percent": psutil.cpu_percent(),
                "memory_percent": psutil.virtual_memory().percent,
                "disk_percent": psutil.disk_usage('/').percent,
                "cpu_count": psutil.cpu_count()
            }
        except ImportError:
            return {"status": "error", "message": "psutil未安装"}
        except Exception as e:
            return {"status": "error", "message": str(e)}


# 创建全局工具服务实例
tool_service = ToolService()
