from celery import Celery
from celery.schedules import crontab
from datetime import timedelta
import os
import subprocess
import threading
from typing import Any, Dict, Optional
from queue import Queue

from src.config.app_config import settings

from src.utils.logger import get_logger

logger = get_logger(__name__)


class LocalTaskQueue:
    """本地任务队列 - 替代Celery的轻量级实现"""
    
    def __init__(self):
        self._queue = Queue()
        self._running = False
        self._worker_thread = None
    
    def start_worker(self):
        """启动本地任务工作线程"""
        if self._running:
            return
        
        self._running = True
        self._worker_thread = threading.Thread(target=self._worker_loop, daemon=True)
        self._worker_thread.start()
    
    def _worker_loop(self):
        """工作线程循环"""
        while self._running:
            try:
                task = self._queue.get(timeout=1)
                func = task['func']
                args = task.get('args', ())
                kwargs = task.get('kwargs', {})
                callback = task.get('callback')
                
                try:
                    result = func(*args, **kwargs)
                    if callback:
                        callback(result)
                except Exception as e:
                    if callback:
                        callback({"status": "error", "message": str(e)})
                finally:
                    self._queue.task_done()
            except Exception as e:
                logger.debug("工作线程任务获取异常", exc_info=True)
    
    def submit(self, func, *args, callback=None, **kwargs):
        """提交任务到队列"""
        self._queue.put({
            'func': func,
            'args': args,
            'kwargs': kwargs,
            'callback': callback
        })
    
    def stop(self):
        """停止工作线程"""
        self._running = False
        if self._worker_thread:
            self._worker_thread.join(timeout=5)


class CeleryService:
    """Celery分布式任务队列服务 - 支持本地模式和远程模式"""
    
    _instance = None
    
    def __new__(cls, *args, **kwargs):
        if not cls._instance:
            cls._instance = super().__new__(cls)
        return cls._instance
    
    def __init__(self, use_local: bool = True):
        if hasattr(self, '_initialized') and self._initialized:
            return
        
        self.use_local = use_local
        
        if use_local:
            # 使用本地任务队列
            self.local_queue = LocalTaskQueue()
            self.local_queue.start_worker()
            self.app = None
        else:
            # 使用Celery
            self.app = Celery(
                'tasks',
                broker=settings.celery_broker_url,
                backend=settings.celery_result_backend
            )
            
            self.app.conf.update(
                task_serializer='json',
                accept_content=['json'],
                result_serializer='json',
                timezone='Asia/Shanghai',
                enable_utc=True,
                task_time_limit=300,
                task_soft_time_limit=240,
                worker_concurrency=4,
            )
            
            self._register_tasks()
        
        self._initialized = True
    
    def _register_tasks(self):
        """注册Celery任务"""
        if not self.app:
            return
        
        @self.app.task(name='tasks.cleanup_old_logs')
        def cleanup_old_logs():
            return self._cleanup_old_logs()
        
        @self.app.task(name='tasks.system_health_check')
        def system_health_check():
            return self._system_health_check()
        
        @self.app.task(name='tasks.database_backup')
        def database_backup():
            return self._database_backup()
        
        @self.app.task(name='tasks.process_command')
        def process_command(command: str, timeout: int = 30):
            return self._process_command(command, timeout)
        
        @self.app.task(name='tasks.process_conversation')
        def process_conversation(conversation_data: dict):
            return self._process_conversation(conversation_data)
    
    def submit_task(self, task_name: str, *args, callback=None, **kwargs) -> Any:
        """提交任务（统一接口，自动选择本地或远程）"""
        if self.use_local:
            # 本地模式：直接获取函数并执行
            task_func = self._get_local_task(task_name)
            if task_func:
                def wrapper():
                    try:
                        result = task_func(*args, **kwargs)
                        if callback:
                            callback(result)
                        return result
                    except Exception as e:
                        if callback:
                            callback({"status": "error", "message": str(e)})
                        return {"status": "error", "message": str(e)}
                
                self.local_queue.submit(wrapper)
                return {"status": "queued", "message": "任务已提交到本地队列"}
            else:
                return {"status": "error", "message": f"任务 {task_name} 未找到"}
        else:
            # Celery模式
            async_result = self.app.send_task(task_name, args=args, kwargs=kwargs)
            return {"status": "queued", "task_id": async_result.task_id}
    
    def _get_local_task(self, task_name: str):
        """获取本地任务函数"""
        tasks = {
            'tasks.cleanup_old_logs': self._cleanup_old_logs,
            'tasks.system_health_check': self._system_health_check,
            'tasks.database_backup': self._database_backup,
            'tasks.process_command': self._process_command,
            'tasks.process_conversation': self._process_conversation,
            'tasks.process_audio': self._process_audio,
            'tasks.process_video': self._process_video,
            'tasks.process_document': self._process_document,
            'tasks.generate_report': self._generate_report,
            'tasks.perform_analysis': self._perform_analysis,
        }
        return tasks.get(task_name)
    
    def _cleanup_old_logs(self) -> Dict[str, Any]:
        """清理旧日志文件"""
        import shutil
        from datetime import datetime, timedelta
        
        log_dir = settings.log_dir
        if not os.path.exists(log_dir):
            return {"status": "success", "message": "Log directory not found"}
        
        cutoff_date = datetime.now() - timedelta(days=7)
        cleaned_count = 0
        
        for filename in os.listdir(log_dir):
            filepath = os.path.join(log_dir, filename)
            if os.path.isfile(filepath):
                try:
                    file_time = datetime.fromtimestamp(os.path.getmtime(filepath))
                    if file_time < cutoff_date:
                        os.remove(filepath)
                        cleaned_count += 1
                except Exception as e:
                    logger.debug(f"清理日志文件失败: {filepath}", exc_info=True)
        
        return {"status": "success", "cleaned_count": cleaned_count}
    
    def _system_health_check(self) -> Dict[str, Any]:
        """系统健康检查"""
        from src.services.system_monitor import system_monitor
        
        try:
            sys_info = system_monitor.get_system_info()
            proc_info = system_monitor.get_process_info()
            
            return {
                "status": "success",
                "cpu_percent": sys_info.cpu_percent,
                "memory_percent": sys_info.memory_percent,
                "disk_percent": sys_info.disk_percent,
                "process_cpu": proc_info.cpu_percent,
                "process_memory": proc_info.memory_percent
            }
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def _database_backup(self) -> Dict[str, Any]:
        """数据库备份"""
        import shutil
        from datetime import datetime
        
        db_dir = settings.db_dir
        backup_dir = os.path.join(db_dir, 'backups')
        
        os.makedirs(backup_dir, exist_ok=True)
        
        backup_count = 0
        for filename in os.listdir(db_dir):
            if filename.endswith('.sqlite'):
                src_path = os.path.join(db_dir, filename)
                backup_name = f"{filename}.{datetime.now().strftime('%Y%m%d_%H%M%S')}.bak"
                dst_path = os.path.join(backup_dir, backup_name)
                
                try:
                    shutil.copy2(src_path, dst_path)
                    backup_count += 1
                except Exception as e:
                    logger.error(f"Backup failed for {filename}: {e}", exc_info=True)
        
        return {"status": "success", "backed_up_count": backup_count}
    
    def _process_command(self, command: str, timeout: int = 30) -> Dict[str, Any]:
        """执行命令任务"""
        from src.services.command_executor import command_executor
        
        stdout, stderr, returncode = command_executor.run_simple_command(command, timeout)
        
        return {
            "status": "success" if returncode == 0 else "error",
            "stdout": stdout,
            "stderr": stderr,
            "returncode": returncode
        }
    
    def _process_conversation(self, conversation_data: dict) -> Dict[str, Any]:
        """处理对话数据"""
        from src.services.database_service import db_service
        
        try:
            record_id = db_service.save_conversation(
                conversation_count=conversation_data.get('conversation_count', 1),
                conversation_content=conversation_data.get('conversation_content', ''),
                permission_level=conversation_data.get('permission_level', 'normal'),
                session_id=conversation_data.get('session_id')
            )
            
            return {"status": "success", "record_id": record_id}
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def _process_audio(self, audio_path: str, operation: str = "transcribe", **kwargs) -> Dict[str, Any]:
        """处理音频文件"""
        try:
            if operation == "transcribe":
                import speech_recognition as sr
                r = sr.Recognizer()
                with sr.AudioFile(audio_path) as source:
                    audio = r.record(source)
                text = r.recognize_google(audio, language='zh-CN')
                return {"status": "success", "transcription": text}
            
            elif operation == "extract_features":
                import librosa
                y, sr = librosa.load(audio_path)
                tempo, beat_frames = librosa.beat.beat_track(y=y, sr=sr)
                mfcc = librosa.feature.mfcc(y=y, sr=sr)
                
                return {
                    "status": "success",
                    "tempo": float(tempo),
                    "mfcc_mean": float(mfcc.mean()),
                    "mfcc_std": float(mfcc.std())
                }
            
            elif operation == "convert":
                from pydub import AudioSegment
                audio = AudioSegment.from_file(audio_path)
                output_path = audio_path.replace('.wav', '_converted.mp3')
                audio.export(output_path, format='mp3')
                return {"status": "success", "output_path": output_path}
            
            else:
                return {"status": "error", "message": f"Unknown operation: {operation}"}
        
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def _process_video(self, video_path: str, operation: str = "extract_frames", **kwargs) -> Dict[str, Any]:
        """处理视频文件"""
        try:
            if operation == "extract_frames":
                import decord
                vr = decord.VideoReader(video_path)
                frame_count = len(vr)
                fps = vr.get_avg_fps()
                
                return {
                    "status": "success",
                    "frame_count": frame_count,
                    "fps": float(fps),
                    "duration": frame_count / fps
                }
            
            elif operation == "cut":
                from moviepy.editor import VideoFileClip
                start = kwargs.get('start', 0)
                end = kwargs.get('end', 10)
                
                clip = VideoFileClip(video_path)
                subclip = clip.subclip(start, end)
                output_path = video_path.replace('.mp4', '_cut.mp4')
                subclip.write_videofile(output_path)
                
                return {"status": "success", "output_path": output_path}
            
            elif operation == "extract_audio":
                from moviepy.editor import VideoFileClip
                video = VideoFileClip(video_path)
                audio = video.audio
                output_path = video_path.replace('.mp4', '.mp3')
                audio.write_audiofile(output_path)
                
                return {"status": "success", "output_path": output_path}
            
            else:
                return {"status": "error", "message": f"Unknown operation: {operation}"}
        
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def _process_document(self, doc_path: str, operation: str = "extract_text", **kwargs) -> Dict[str, Any]:
        """处理文档文件"""
        try:
            if operation == "extract_text":
                _, ext = os.path.splitext(doc_path)
                
                if ext.lower() == '.pdf':
                    import pdfplumber
                    with pdfplumber.open(doc_path) as pdf:
                        text = '\n'.join(page.extract_text() for page in pdf.pages)
                
                elif ext.lower() == '.docx':
                    from docx import Document
                    doc = Document(doc_path)
                    text = '\n'.join(para.text for para in doc.paragraphs)
                
                elif ext.lower() == '.xlsx':
                    import openpyxl
                    wb = openpyxl.load_workbook(doc_path)
                    text = ""
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        for row in ws.iter_rows(values_only=True):
                            text += ' '.join(str(cell) for cell in row if cell) + '\n'
                
                elif ext.lower() == '.pptx':
                    from pptx import Presentation
                    prs = Presentation(doc_path)
                    text = ""
                    for slide in prs.slides:
                        for shape in slide.shapes:
                            if hasattr(shape, 'text'):
                                text += shape.text + '\n'
                
                else:
                    return {"status": "error", "message": f"Unsupported file type: {ext}"}
                
                return {"status": "success", "text": text[:5000]}
            
            elif operation == "convert_to_text":
                _, ext = os.path.splitext(doc_path)
                
                if ext.lower() == '.pdf':
                    import pdfplumber
                    with pdfplumber.open(doc_path) as pdf:
                        text = '\n'.join(page.extract_text() for page in pdf.pages)
                elif ext.lower() == '.docx':
                    from docx import Document
                    doc = Document(doc_path)
                    text = '\n'.join(para.text for para in doc.paragraphs)
                elif ext.lower() == '.xlsx':
                    import openpyxl
                    wb = openpyxl.load_workbook(doc_path)
                    text = ""
                    for sheet in wb.sheetnames:
                        ws = wb[sheet]
                        for row in ws.iter_rows(values_only=True):
                            text += ' '.join(str(cell) for cell in row if cell) + '\n'
                else:
                    return {"status": "error", "message": f"Unsupported file type: {ext}"}
                
                output_path = doc_path.replace(ext, '.txt')
                with open(output_path, 'w', encoding='utf-8') as f:
                    f.write(text)
                
                return {"status": "success", "output_path": output_path}
            
            else:
                return {"status": "error", "message": f"Unknown operation: {operation}"}
        
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def _generate_report(self, data: dict, format: str = "text") -> Dict[str, Any]:
        """生成报告"""
        try:
            if format == "text":
                report = f"报告生成时间: {datetime.now().isoformat()}\n"
                report += "=" * 50 + "\n"
                for key, value in data.items():
                    report += f"{key}: {value}\n"
                return {"status": "success", "report": report}
            
            elif format == "docx":
                from docx import Document
                doc = Document()
                doc.add_heading('AI分析报告', level=1)
                
                for key, value in data.items():
                    doc.add_heading(key, level=2)
                    doc.add_paragraph(str(value))
                
                output_path = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.docx"
                doc.save(output_path)
                
                return {"status": "success", "output_path": output_path}
            
            elif format == "pptx":
                from pptx import Presentation
                prs = Presentation()
                
                slide = prs.slides.add_slide(prs.slide_layouts[0])
                title = slide.shapes.title
                title.text = "AI分析报告"
                
                for key, value in data.items():
                    slide = prs.slides.add_slide(prs.slide_layouts[1])
                    slide.shapes.title.text = key
                    slide.placeholders[1].text = str(value)
                
                output_path = f"report_{datetime.now().strftime('%Y%m%d_%H%M%S')}.pptx"
                prs.save(output_path)
                
                return {"status": "success", "output_path": output_path}
            
            else:
                return {"status": "error", "message": f"Unknown format: {format}"}
        
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def _perform_analysis(self, data: dict, analysis_type: str = "basic") -> Dict[str, Any]:
        """执行数据分析"""
        try:
            if analysis_type == "basic":
                return {
                    "status": "success",
                    "analysis": {
                        "data_length": len(str(data)),
                        "keys_count": len(data.keys()) if isinstance(data, dict) else 0,
                        "analysis_time": datetime.now().isoformat()
                    }
                }
            
            elif analysis_type == "financial":
                import pandas as pd
                import pandas_ta as ta
                
                df = pd.DataFrame(data.get('data', []))
                if 'close' in df.columns:
                    df.ta.rsi(append=True)
                    df.ta.macd(append=True)
                    
                    return {
                        "status": "success",
                        "indicators": {
                            "rsi": df['RSI_14'].tail(10).tolist(),
                            "macd": df['MACD_12_26_9'].tail(10).tolist()
                        }
                    }
                else:
                    return {"status": "error", "message": "No 'close' column found"}
            
            elif analysis_type == "text":
                from gensim.summarization import summarize
                text = data.get('text', '')
                summary = summarize(text, ratio=0.2)
                
                return {"status": "success", "summary": summary}
            
            else:
                return {"status": "error", "message": f"Unknown analysis type: {analysis_type}"}
        
        except Exception as e:
            return {"status": "error", "message": str(e)}
    
    def close(self):
        """关闭服务"""
        if self.use_local and self.local_queue:
            self.local_queue.stop()


# 初始化Celery服务（使用本地模式）
celery_service = CeleryService(use_local=True)
